#!/usr/bin/env python3
"""
完善答辩演示PPT - 使用python-pptx
结合项目简介PPT的内容进行增强
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# 文件路径
INPUT_PPTX = "D:/ASUS/Documents/jiedan/waibao/waibu/计算机设计大赛答辩演示_无敌暴龙战队.pptx"
OUTPUT_PPTX = "D:/ASUS/Documents/jiedan/waibao/waibu/计算机设计大赛答辩演示_无敌暴龙战队_完善版.pptx"

def add_team_slide(prs):
    """添加团队介绍页（在第13页之前）"""
    # 使用第一个可用的布局（通常是空白或标题页）
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    
    # 添加标题背景
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    title_shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "团队介绍"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title_frame.paragraphs[0].font.name = "微软雅黑"
    
    # 团队成员信息
    team_members = [
        ("侯宗阳", "队长 / 开发", "系统架构设计与核心算法实现\nJava / Spring Boot / 数据库"),
        ("王晨溪", "前端开发", "React组件开发与可视化实现\nReact / ECharts / WebSocket"),
        ("徐永佳", "算法工程师", "威胁检测算法设计与优化\nPython / 数据分析 / 机器学习"),
        ("徐慧芳", "测试运维", "Docker部署与系统测试\nDocker / Linux / 性能调优"),
        ("王昭烨", "后端开发", "后端服务开发与数据处理\nJava / MySQL / RESTful API"),
    ]
    
    # 添加团队成员卡片（5人布局：2+2+1）
    for i, (name, role, desc) in enumerate(team_members):
        if i < 4:
            # 前4人：2x2布局
            col = i % 2
            row = i // 2
            x = Inches(0.5 + col * 4.8)
            y = Inches(1.2 + row * 1.7)
        else:
            # 第5人：居中放在底部
            x = Inches(2.85)  # (10 - 4.3) / 2 = 2.85 居中
            y = Inches(4.6)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.3), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        
        # 姓名和角色
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), Inches(3.9), Inches(0.35))
        name_frame = name_box.text_frame
        name_frame.text = f"{name}  |  {role}"
        name_frame.paragraphs[0].font.size = Pt(15)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = RGBColor(0x0F, 0x4C, 0x75)
        name_frame.paragraphs[0].font.name = "微软雅黑"
        
        # 描述
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.5), Inches(3.9), Inches(0.9))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(11)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        desc_frame.paragraphs[0].font.name = "微软雅黑"
    
    # 添加参赛编号
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.4))
    team_frame = team_box.text_frame
    team_frame.text = "无敌暴龙战队  |  参赛编号：T2601487"
    team_frame.paragraphs[0].font.size = Pt(14)
    team_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    team_frame.paragraphs[0].font.name = "微软雅黑"
    
    return slide

def enhance_deployment_slide(slide):
    """完善部署页面，添加Swagger地址"""
    # 查找文本框并添加API文档地址
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "后端API:" in text and "swagger" not in text.lower():
                # 在现有文本后添加API文档信息
                new_text = text.replace(
                    "后端API:",
                    "API文档:\nhttp://8.146.228.64:8080/swagger-ui\n\n后端API:"
                )
                shape.text_frame.text = new_text
                # 设置字体
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = "微软雅黑"

def main():
    print("=" * 60)
    print("完善答辩演示PPT")
    print("=" * 60)
    
    # 加载PPT
    prs = Presentation(INPUT_PPTX)
    print(f"[OK] 已加载PPT，共 {len(prs.slides)} 页")
    
    # 1. 在第12页（部署页）后添加团队介绍页
    print("\n添加团队介绍页...")
    team_slide = add_team_slide(prs)
    
    # 2. 完善部署页面
    print("完善部署页面...")
    # 找到部署页面（第12页）
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "部署与运行" in shape.text_frame.text:
                print(f"  找到部署页面（第{i+1}页）")
                enhance_deployment_slide(slide)
                break
    
    # 3. 保存PPT
    print("\n保存完善版PPT...")
    prs.save(OUTPUT_PPTX)
    
    print("\n" + "=" * 60)
    print(f"[OK] 完善版PPT已生成: {OUTPUT_PPTX}")
    print(f"[OK] 共 {len(prs.slides)} 页")
    print("=" * 60)
    print("\n主要改进:")
    print("1. 新增团队介绍页（5位成员详细信息）")
    print("2. 完善了部署页面，添加Swagger API文档地址")
    print("3. 保持了答辩PPT的精炼风格")

if __name__ == "__main__":
    main()
